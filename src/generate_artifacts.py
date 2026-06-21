from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.platypus import (
    Image,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)


def project_root() -> Path:
    return Path(__file__).resolve().parents[1]


def register_fonts() -> str:
    arial = Path("C:/Windows/Fonts/arial.ttf")
    if arial.exists():
        pdfmetrics.registerFont(TTFont("Arial", str(arial)))
        return "Arial"
    return "Helvetica"


def table_from_dataframe(df: pd.DataFrame, max_rows: int = 10) -> Table:
    shown = df.head(max_rows).copy()
    data = [shown.columns.tolist()] + shown.round(4).astype(str).values.tolist()
    table = Table(data, repeatRows=1)
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1F4E79")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("GRID", (0, 0), (-1, -1), 0.25, colors.grey),
                ("FONTNAME", (0, 0), (-1, -1), "Arial"),
                ("FONTSIZE", (0, 0), (-1, -1), 8),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F2F6FA")]),
            ]
        )
    )
    return table


def add_report_image(story: list, path: Path, width: float = 6.4 * inch) -> None:
    if path.exists():
        story.append(Image(str(path), width=width, height=width * 0.58))
        story.append(Spacer(1, 0.16 * inch))


def generate_pdf(root: Path) -> Path:
    font = register_fonts()
    output_path = root / "reports" / "final_report.pdf"
    table_dir = root / "reports" / "tables"
    figure_dir = root / "reports" / "figures"

    styles = getSampleStyleSheet()
    styles.add(
        ParagraphStyle(
            name="VNTitle",
            parent=styles["Title"],
            fontName=font,
            fontSize=18,
            leading=24,
            spaceAfter=16,
        )
    )
    styles.add(
        ParagraphStyle(
            name="VNHeading",
            parent=styles["Heading2"],
            fontName=font,
            fontSize=13,
            leading=16,
            spaceBefore=12,
            spaceAfter=8,
        )
    )
    styles.add(
        ParagraphStyle(
            name="VNBody",
            parent=styles["BodyText"],
            fontName=font,
            fontSize=10,
            leading=14,
            spaceAfter=6,
        )
    )

    metrics = pd.read_csv(table_dir / "model_metrics.csv")
    metrics_display = metrics[
        ["model", "threshold", "precision", "recall", "f1_score", "roc_auc", "pr_auc"]
    ]
    spark_summary = pd.read_csv(table_dir / "spark_dataset_summary.csv")
    category = pd.read_csv(table_dir / "spark_fraud_by_category.csv")
    merchants = pd.read_csv(table_dir / "spark_high_risk_merchants.csv")

    story = [
        Paragraph("Fraud Detection System using Big Data Analytics", styles["VNTitle"]),
        Paragraph(
            "Báo cáo tổng hợp pipeline phát hiện giao dịch gian lận bằng EDA, preprocessing, Apache Spark và Machine Learning.",
            styles["VNBody"],
        ),
        Paragraph("1. Dataset", styles["VNHeading"]),
        Paragraph(
            "Dataset gồm fraudTrain.csv và fraudTest.csv từ Kaggle. Train set có 1,296,675 giao dịch, trong đó 7,506 giao dịch là fraud. Fraud rate chỉ khoảng 0.5789%, thể hiện bài toán mất cân bằng mạnh.",
            styles["VNBody"],
        ),
        table_from_dataframe(spark_summary),
        Paragraph("2. Liên hệ Big Data 5V", styles["VNHeading"]),
        Paragraph(
            "Volume: hơn 1.85 triệu giao dịch train + test. Velocity: giao dịch tài chính trong thực tế phát sinh liên tục và cần cảnh báo nhanh. Variety: dữ liệu gồm amount, time, location, merchant, category và hồ sơ khách hàng. Veracity: pipeline kiểm tra missing/duplicate và loại bỏ cột định danh trực tiếp. Value: mô hình hỗ trợ phát hiện fraud, giảm rủi ro tài chính và hỗ trợ ra quyết định.",
            styles["VNBody"],
        ),
        Paragraph(
            "Dự án bao gồm ba nhiệm vụ chính của dữ liệu lớn: Data Management, Data Modeling and Analytics, Visualization/Decision/Value.",
            styles["VNBody"],
        ),
        Paragraph("3. EDA", styles["VNHeading"]),
        Paragraph(
            "EDA phân tích target distribution, amount, category, merchant, thời gian, vị trí địa lý và khoảng cách customer-merchant.",
            styles["VNBody"],
        ),
    ]
    add_report_image(story, figure_dir / "target_distribution.png")
    add_report_image(story, figure_dir / "category_transaction_count_and_fraud_rate.png")

    story.extend(
        [
            PageBreak(),
            Paragraph("4. Apache Spark Processing", styles["VNHeading"]),
            Paragraph(
                "Spark DataFrame được dùng để tính fraud theo category, state, merchant và hour. Kết quả được xuất ra reports/tables.",
                styles["VNBody"],
            ),
            table_from_dataframe(category, max_rows=8),
            Spacer(1, 0.18 * inch),
            Paragraph("Top merchant có tỷ lệ fraud cao:", styles["VNBody"]),
            table_from_dataframe(merchants, max_rows=8),
            Paragraph("5. Modeling", styles["VNHeading"]),
            Paragraph(
                "Ba mô hình được huấn luyện: Logistic Regression, Random Forest và HistGradientBoosting. Training dùng sampling giữ toàn bộ fraud và lấy non-fraud theo tỷ lệ 5:1, tune threshold theo F1-score trên validation split, sau đó đánh giá trên toàn bộ test set.",
                styles["VNBody"],
            ),
            table_from_dataframe(metrics_display),
        ]
    )
    add_report_image(story, figure_dir / "roc_curve_models.png")
    add_report_image(story, figure_dir / "precision_recall_curve_models.png")
    add_report_image(story, figure_dir / "model_metrics_comparison.png")

    story.extend(
        [
            PageBreak(),
            Paragraph("6. Kết luận", styles["VNHeading"]),
            Paragraph(
                "HistGradientBoosting tuned F1 là mô hình tốt nhất trong thử nghiệm, đạt Accuracy khoảng 0.9932, Precision khoảng 0.3557, Recall khoảng 0.9245 và F1-score khoảng 0.5137. ROC-AUC đạt khoảng 0.9976 và PR-AUC đạt khoảng 0.8701. So với Random Forest, mô hình này cải thiện rõ rệt cả khả năng phát hiện fraud và chất lượng cảnh báo.",
                styles["VNBody"],
            ),
            Paragraph("7. Hướng phát triển", styles["VNHeading"]),
            Paragraph(
                "Có thể cải thiện bằng threshold tuning, PR-AUC, XGBoost/LightGBM, target encoding có regularization và triển khai batch scoring bằng Spark.",
                styles["VNBody"],
            ),
        ]
    )

    doc = SimpleDocTemplate(str(output_path), pagesize=A4, rightMargin=36, leftMargin=36, topMargin=36, bottomMargin=36)
    doc.build(story)
    return output_path


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullets_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for idx, bullet in enumerate(bullets):
        paragraph = body.paragraphs[0] if idx == 0 else body.add_paragraph()
        paragraph.text = bullet
        paragraph.font.size = Pt(22)
        paragraph.level = 0


def add_picture_slide(prs: Presentation, title: str, image_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(0.8), Inches(1.35), width=Inches(8.4))


def add_metrics_slide(prs: Presentation, metrics: pd.DataFrame) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Model Evaluation"
    rows, cols = metrics.shape[0] + 1, metrics.shape[1]
    table = slide.shapes.add_table(rows, cols, Inches(0.35), Inches(1.25), Inches(9.3), Inches(1.8)).table

    for col_idx, col in enumerate(metrics.columns):
        table.cell(0, col_idx).text = col

    for row_idx, (_, row) in enumerate(metrics.iterrows(), start=1):
        for col_idx, col in enumerate(metrics.columns):
            value = row[col]
            table.cell(row_idx, col_idx).text = f"{value:.4f}" if isinstance(value, float) else str(value)


def generate_pptx(root: Path) -> Path:
    output_path = root / "reports" / "fraud_detection_presentation.pptx"
    table_dir = root / "reports" / "tables"
    figure_dir = root / "reports" / "figures"
    metrics = pd.read_csv(table_dir / "model_metrics.csv")[
        ["model", "threshold", "precision", "recall", "f1_score", "roc_auc", "pr_auc"]
    ]

    prs = Presentation()
    add_title_slide(
        prs,
        "Fraud Detection System",
        "Big Data Analytics with PySpark and Machine Learning",
    )
    add_bullets_slide(
        prs,
        "Problem Statement",
        [
            "Detect fraudulent financial transactions from historical data.",
            "Dataset is highly imbalanced: fraud rate below 1%.",
            "Need EDA, preprocessing, Spark processing and ML evaluation.",
        ],
    )
    add_bullets_slide(
        prs,
        "Dataset",
        [
            "Source: Kaggle Credit Card Transactions Fraud Detection.",
            "Train: 1,296,675 transactions.",
            "Test: 555,719 transactions.",
            "Target: is_fraud.",
        ],
    )
    add_bullets_slide(
        prs,
        "Big Data 5V Mapping",
        [
            "Volume: more than 1.85M transactions across train and test.",
            "Velocity: financial transactions require fast fraud alerts.",
            "Variety: amount, time, location, merchant, category and customer profile.",
            "Veracity: missing/duplicate checks and direct identifier removal.",
            "Value: fraud detection supports risk reduction and decision making.",
        ],
    )
    add_picture_slide(prs, "Target Distribution", figure_dir / "target_distribution.png")
    add_picture_slide(prs, "Category Fraud Analysis", figure_dir / "category_transaction_count_and_fraud_rate.png")
    add_bullets_slide(
        prs,
            "Preprocessing",
            [
                "Time features: hour, day of week, month, weekend flag.",
                "Customer age and customer-merchant distance.",
                "One-hot encoding for category/gender/state.",
                "Frequency encoding for merchant/job/city.",
                "Numeric scaling and removal of direct identifiers.",
            ],
    )
    add_bullets_slide(
        prs,
        "Apache Spark",
        [
            "Read raw train data using Spark DataFrame.",
            "Aggregate fraud by category, state, merchant and hour.",
            "Spark version: 3.5.1.",
            "Pandas vs Spark category aggregation: Spark was faster on this machine.",
        ],
    )
    add_metrics_slide(prs, metrics)
    add_picture_slide(prs, "ROC Curve", figure_dir / "roc_curve_models.png")
    add_picture_slide(prs, "Precision-Recall Curve", figure_dir / "precision_recall_curve_models.png")
    add_picture_slide(prs, "Confusion Matrix - Random Forest", figure_dir / "confusion_matrix_random_forest.png")
    add_bullets_slide(
        prs,
            "Conclusion",
            [
                "HistGradientBoosting achieved the best overall performance.",
                "Recall: 0.9245, F1-score: 0.5137.",
                "ROC-AUC: 0.9976, PR-AUC: 0.8701.",
                "F1-score improved from Random Forest's 0.3011 to 0.5137.",
                "Future work: threshold tuning, PR-AUC, XGBoost/LightGBM.",
            ],
    )

    prs.save(output_path)
    return output_path


def main() -> None:
    root = project_root()
    pdf_path = generate_pdf(root)
    pptx_path = generate_pptx(root)
    print(f"PDF generated: {pdf_path}")
    print(f"PPTX generated: {pptx_path}")


if __name__ == "__main__":
    main()
